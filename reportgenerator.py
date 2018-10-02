#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Created on Fri Jul  7 12:10:30 2017

@author: peter
"""

# Imports
import os # file i/o operations
import glob # file i/o operations
from datetime import date# finding today's date
from pptx import Presentation # creating the PPT
from pptx.util import Inches, Pt
#import comtypes.client # for opening PowerPoint from python

"""
def PPTtoPDF(inputFileName, outputFileName, formatType = 32):
    
    Converts a PPT file to a PDF by opening PowerPoint, opening the file, and 
    then saving as a PowerPoint. Requires Windows. 
    
    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
    powerpoint.Visible = 1

    if outputFileName[-3:] != 'pdf':
        outputFileName = outputFileName + ".pdf"
    deck = powerpoint.Presentations.Open(inputFileName)
    deck.SaveAs(outputFileName, formatType) # formatType = 32 for ppt to pdf
    deck.Close()
    powerpoint.Quit()
"""

def addTitleSlide(title,subtitle):
    slide = prs.slides.add_slide(title_only_slide_layout)
    #left = Inches(2)
    #top = width = height = Inches(1)
    #txBox = slide.shapes.add_textbox(left, top, width, height)
    #tf = txBox.text_frame
    shapes = slide.shapes
    shapes.title.text = title
    #p = tf.add_paragraph()
    #p.text = title
    #p.font.size = Pt(40)

def addImageSlide(imageFileName):
    """
    Adds a full-screen image to a blank full-screen slide.
    """
    slide = prs.slides.add_slide(blank_slide_layout)
    slide.shapes.add_picture(imageFileName, 0, 0, height=prs.slide_height, width=prs.slide_width)

# Get today's date, formatted to MATLAB's default (e.g. 2017-Jul-09)
today = date.today().strftime('%d-%b-%Y')

# make filename
reportFile = today + '_report.pptx'

# Initialize presentation
prs = Presentation()
prs.slide_height = 5143500 # Widescreen aspect ratio
title_slide_layout = prs.slide_layouts[1]
blank_slide_layout = prs.slide_layouts[6] # blank slide

# Create title slide
addTitleSlide('Current Cycling Progress', today)

# CD to directory with most recent images
# os.chdir('C:\\Users\\Arbin\\Box Sync\\Batch images\\' + today + '\\')

# Add .png files in this directory. Start with summary figures
all_images = glob.glob('*.png')
for file in all_images:
    if "summary" in file:
        addImageSlide(file)
        
# Cell "spec sheets"
for file in all_images:
    if "summary" not in file:
        addImageSlide(file)

"""
# Directory for saving reports
saveDir = 'C:\\Users\\Arbin\\Box Sync\\Reports'
os.chdir(saveDir)

# Create file names
reportFileFull = saveDir + '\\' + reportFile
reportFileFullPDF = saveDir + '\\' + reportFile.replace('pptx','pdf')

# Save powerpoint
prs.save(reportFileFull)
# Convert to PDF
#PPTtoPDF(reportFileFull,reportFileFullPDF)
"""